import pandas as pd
import numpy as np
import json  # Импортируем библиотеку json
import matplotlib.pyplot as plt
from pptx import Presentation
from openpyxl import Workbook
import os


# Функция для загрузки данных из файлов
def load_data():
    files = {
        'client': r'C:\Users\vanya\Desktop\ХАХА\client.parquet',
        'physical': r'C:\Users\vanya\Desktop\ХАХА\physical.parquet',
        'company': r'C:\Users\vanya\Desktop\ХАХА\company.parquet',
        'plan': r'C:\Users\vanya\Desktop\ХАХА\plan.json',
        'subscribers': r'C:\Users\vanya\Desktop\ХАХА\subscribers.csv',
        'psxstats': r'C:\Users\vanya\Desktop\ХАХА\psxstats.csv',  # Путь к файлу CSV
        'psxattrs': r'C:\Users\vanya\Desktop\ХАХА\psxattrs.csv'
    }

    data = {}

    for key, filename in files.items():
        if not os.path.exists(filename):
            raise FileNotFoundError(f"Файл '{filename}' не найден.")

        if filename.endswith('.parquet'):
            data[key] = pd.read_parquet(filename)
        elif filename.endswith('.json'):
            with open(filename, 'r', encoding='utf-8') as f:
                data[key] = pd.json_normalize(json.load(f))
        elif filename.endswith('.csv'):
            data[key] = pd.read_csv(filename, encoding='windows-1251')  # Попробуйте указать кодировку

    return (data['client'], data['physical'], data['company'],
            data['plan'], data['subscribers'], data['psxstats'],
            data['psxattrs'])


# Функция для создания и заполнения CSV файла psxstats с 10,000 записями
def create_psxstats_csv():
    # Генерация данных для 10,000 записей
    num_records = 10000

    # Создание списка для хранения данных
    example_data = []

    for i in range(num_records):
        example_data.append([
            123450 + i,  # subscriber_number
            i % 10 + 1,  # client_id (1-10)
            f"2024-11-24 {i // 2:02d}:00:00",  # session_start
            f"2024-11-24 {i // 2 + 1:02d}:00:00",  # session_end
            np.random.randint(1000, 20000),  # uploaded_traffic
            np.random.randint(1000, 20000),  # downloaded_traffic
            i + 1,  # session_number (1-10000)
            f"A{i % 5 + 1}"  # switch_id (A1-A5)
        ])

    # Создание DataFrame и сохранение в CSV файл
    columns = ['subscriber_number', 'client_id', 'session_start',
               'session_end', 'uploaded_traffic',
               'downloaded_traffic', 'session_number', 'switch_id']

    psx_stats_df = pd.DataFrame(example_data, columns=columns)
    psx_stats_df.to_csv(r'C:\Users\vanya\Desktop\ХАХА\psxstats.csv', index=False)


# Функция для анализа трафика и подсчета взломов
def analyze_traffic(psx_stats):
    if 'subscriber_number' not in psx_stats.columns:
        raise ValueError("Столбец 'subscriber_number' отсутствует в psx_stats.")

    # Случайное количество взломов для примера (можно заменить на реальные данные)
    hacked_count = np.random.randint(0, 1000)
    return hacked_count


# Функция для контроля качества данных
def data_quality_report(data):
    return {
        'total_rows': len(data),
        'missing_values': data.isnull().sum().to_dict(),
        'duplicates': data.duplicated().sum()
    }


# Функция для сохранения отчетов в Excel
def save_reports_to_excel(hourly_reports):
    wb = Workbook()

    report_sheet = wb.active
    report_sheet.title = "Hourly Reports"

    report_headers = ['hour',
                      'client_id',
                      'uploaded_traffic',
                      'downloaded_traffic',
                      'session_number',
                      'switch_id',
                      'subscriber_number',
                      'start_time',
                      'end_time',
                      'status',
                      'justification']

    report_sheet.append(report_headers)

    for hour in hourly_reports:
        report_sheet.append([hour])  # Добавьте данные отчета

    wb.save(r'C:\Users\vanya\Desktop\ХАХА\hourly_reports.xlsx')


# Функция для создания диаграммы распределения трафика и взломов
def create_traffic_distribution_chart(data):
    plt.figure(figsize=(10, 6))

    # График трафика
    plt.hist(data['uploaded_traffic'].dropna(), bins=50, alpha=0.5, label='Uploaded Traffic')
    plt.hist(data['downloaded_traffic'].dropna(), bins=50, alpha=0.5, label='Downloaded Traffic')

    plt.title('Распределение трафика абонентов')
    plt.xlabel('Объем трафика')
    plt.ylabel('Количество')

    plt.legend()

    plt.savefig('traffic_distribution.png')  # Сохраняем диаграмму

    plt.figure(figsize=(10, 6))

    # График взломов
    hacked_count = np.random.randint(0, 1000)  # Случайное количество взломов

    plt.bar(['Hacked'], [hacked_count], color='red')

    plt.title('Количество взломов')
    plt.ylabel('Количество')

    plt.grid(axis='y')  # Добавляем сетку по оси Y для лучшей читаемости

    plt.savefig('hacked_distribution.png')  # Сохраняем диаграмму взломов

    plt.show()


# Функция для создания презентации
def create_presentation(hourly_reports, quality_report):
    prs = Presentation()

    slide_title = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide_title.shapes.title
    title.text = "Отчет по трафику"

    slide_quality = prs.slides.add_slide(prs.slide_layouts[1])
    quality_title = slide_quality.shapes.title
    quality_title.text = "Качество данных"

    content = slide_quality.shapes.placeholders[1]
    content.text = (f"Всего строк: {quality_report['total_rows']}\n"
                    f"Недостающие значения: {quality_report['missing_values']}\n"
                    f"Дубликаты: {quality_report['duplicates']}")

    prs.save('traffic_report.pptx')


# Основной блок выполнения программы
if __name__ == "__main__":
    try:
        create_psxstats_csv()  # Создание и заполнение CSV файла psxstats

        clients, physical, companies, plans, subscribers, psx_stats, psx_attrs = load_data()

        required_columns = ['subscriber_number',
                            'client_id',
                            'session_start',
                            'session_end',
                            'uploaded_traffic',
                            'downloaded_traffic',
                            'session_number',
                            'switch_id']

        for col in required_columns:
            if col not in psx_stats.columns:
                print(f"Столбец '{col}' отсутствует в psx_stats.")

        hacked_count = analyze_traffic(psx_stats)  # Получаем количество взломов

        quality_report = data_quality_report(psx_stats)

        save_reports_to_excel({})  # Передаем пустой словарь или реализуйте логику отчетов

        create_traffic_distribution_chart(psx_stats)  # Создаем графики

        create_presentation({}, quality_report)  # Передаем пустой словарь или реализуйте логику отчетов

        print("Процесс завершен успешно.")

    except Exception as e:
        print(f"Произошла ошибка: {e}")